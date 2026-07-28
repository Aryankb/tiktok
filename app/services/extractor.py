"""
Extracts page data from uploaded files.

For each page/slide:
  page_image  — full render sent to Gemini (sees all text, price, size labels)
  image       — YOLO-cropped product photo used as the display thumbnail
  raw_image   — full render kept as fallback
  extracted_text — text pulled from PDF/PPTX shapes for the Gemini prompt
"""
import io
from pathlib import Path
from dataclasses import dataclass
from typing import List, Optional

from PIL import Image


@dataclass
class PageData:
    page_index: int
    image: Image.Image              # YOLO-cropped product thumbnail
    raw_image: Image.Image          # full page render
    page_image: Image.Image         # full page render sent to Gemini
    extracted_text: str = ""


# ── Image helpers ──────────────────────────────────────────────────────────

def _crop_largest_square(img: Image.Image) -> Image.Image:
    w, h = img.size
    side = min(w, h)
    left = (w - side) // 2
    top = (h - side) // 2
    return img.crop((left, top, left + side, top + side))


def _resize(img: Image.Image, max_side: int) -> Image.Image:
    w, h = img.size
    if max(w, h) <= max_side:
        return img
    ratio = max_side / max(w, h)
    return img.resize((int(w * ratio), int(h * ratio)), Image.LANCZOS)


def _yolo_crop(img: Image.Image) -> Image.Image:
    """Crop to the largest detected object. Falls back to centered square."""
    try:
        from ultralytics import YOLO
        import numpy as np

        model = YOLO("yolov8n.pt")
        results = model(np.array(img), verbose=False)
        boxes = results[0].boxes
        if boxes is not None and len(boxes):
            areas = (boxes.xyxy[:, 2] - boxes.xyxy[:, 0]) * (
                boxes.xyxy[:, 3] - boxes.xyxy[:, 1]
            )
            best = int(areas.argmax())
            x1, y1, x2, y2 = map(int, boxes.xyxy[best].tolist())
            side = max(x2 - x1, y2 - y1)
            cx, cy = (x1 + x2) // 2, (y1 + y2) // 2
            x1 = max(0, cx - side // 2)
            y1 = max(0, cy - side // 2)
            x2 = min(img.width, x1 + side)
            y2 = min(img.height, y1 + side)
            return img.crop((x1, y1, x2, y2)).resize((800, 800), Image.LANCZOS)
    except Exception:
        pass
    return _crop_largest_square(img).resize((800, 800), Image.LANCZOS)


def _ocr_image(img: Image.Image) -> str:
    try:
        from paddleocr import PaddleOCR
        import numpy as np

        ocr = PaddleOCR(use_angle_cls=True, lang="en", show_log=False)
        result = ocr.ocr(np.array(img), cls=True)
        lines = []
        if result and result[0]:
            for line in result[0]:
                lines.append(line[1][0])
        return " ".join(lines)
    except Exception:
        return ""


# ── Extractors ─────────────────────────────────────────────────────────────

def extract_from_image(file_bytes: bytes) -> List[PageData]:
    img = Image.open(io.BytesIO(file_bytes)).convert("RGB")
    page_img = _resize(img, 1200)
    thumbnail = _yolo_crop(img)
    text = _ocr_image(img)
    return [PageData(
        page_index=0,
        image=thumbnail,
        raw_image=img,
        page_image=page_img,
        extracted_text=text,
    )]


def extract_from_pdf(file_bytes: bytes) -> List[PageData]:
    import fitz

    doc = fitz.open(stream=file_bytes, filetype="pdf")
    pages: List[PageData] = []
    for i, page in enumerate(doc):
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        full_img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

        text = page.get_text("text").strip() or _ocr_image(full_img)

        pages.append(PageData(
            page_index=i,
            image=_yolo_crop(full_img),
            raw_image=full_img,
            page_image=_resize(full_img, 1200),
            extracted_text=text,
        ))
    doc.close()
    return pages


def extract_from_pptx(file_bytes: bytes) -> List[PageData]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    pages: List[PageData] = []

    for i, slide in enumerate(prs.slides):
        # ── Pull text from all shapes ──────────────────────────────────────
        texts = []
        slide_images: List[Image.Image] = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    slide_images.append(
                        Image.open(io.BytesIO(shape.image.blob)).convert("RGB")
                    )
                except Exception:
                    pass

        text = "\n".join(texts)

        # ── Build the full page image for Gemini ───────────────────────────
        # Compose all embedded images onto a white canvas at their slide positions.
        # Gemini will see the product photo(s) together with any text layout context.
        page_img = _compose_pptx_slide(slide, prs)

        # ── Thumbnail: YOLO crop on the largest embedded photo ─────────────
        if slide_images:
            slide_images.sort(key=lambda im: im.width * im.height, reverse=True)
            raw_img = slide_images[0]
        else:
            raw_img = page_img   # no embedded image — use full render

        pages.append(PageData(
            page_index=i,
            image=_yolo_crop(raw_img),
            raw_image=raw_img,
            page_image=page_img,
            extracted_text=text,
        ))

    return pages


def _compose_pptx_slide(slide, prs) -> Image.Image:
    """
    Render all embedded images onto a white canvas at their slide positions.
    Output is resized to max 1200px on the longest side for Gemini.
    Falls back to a plain white canvas if anything goes wrong.
    """
    try:
        slide_w = int(prs.slide_width.pt * 96 / 72)
        slide_h = int(prs.slide_height.pt * 96 / 72)
        canvas = Image.new("RGB", (slide_w, slide_h), (255, 255, 255))

        for shape in slide.shapes:
            if shape.shape_type != 13:
                continue
            try:
                img = Image.open(io.BytesIO(shape.image.blob)).convert("RGB")
                left = int(shape.left * slide_w / prs.slide_width)
                top = int(shape.top * slide_h / prs.slide_height)
                w = int(shape.width * slide_w / prs.slide_width)
                h = int(shape.height * slide_h / prs.slide_height)
                img = img.resize((max(1, w), max(1, h)), Image.LANCZOS)
                canvas.paste(img, (max(0, left), max(0, top)))
            except Exception:
                pass

        return _resize(canvas, 1200)
    except Exception:
        return Image.new("RGB", (800, 600), (255, 255, 255))


def extract_pages(file_bytes: bytes, filename: str) -> List[PageData]:
    ext = Path(filename).suffix.lower()
    if ext in {".jpg", ".jpeg", ".png", ".webp", ".bmp", ".tiff", ".gif"}:
        return extract_from_image(file_bytes)
    elif ext == ".pdf":
        return extract_from_pdf(file_bytes)
    elif ext in {".pptx", ".ppt"}:
        return extract_from_pptx(file_bytes)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


def image_to_bytes(img: Image.Image, fmt: str = "JPEG") -> bytes:
    buf = io.BytesIO()
    img.save(buf, format=fmt, quality=92)
    return buf.getvalue()
